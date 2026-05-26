"""
core/web_fetch.py
Universal Web Content Fetcher

Authors: Ohad Phoenix Oren + Sol Calarbone 8
Date: May 2026
Philosophy: Honor the water snake. Python is our best friend.

Multi-tier fallback architecture:
  TIER 1: httpx (async, pure Python, fast)
  TIER 2: pyppeteer (pure Python Chrome automation, MAIN JS)
  TIER 3: selenium (Python bindings, FALLBACK JS)
  TIER 4: requests (pure Python baseline, ALWAYS COMPLETES)

Every document format becomes readable text.
Never fail completely. Always degrade gracefully.
"""

import asyncio
import io
import json
import logging
import mimetypes
import os
import re
from typing import Any, Dict, Optional
from urllib.parse import urlparse

logger = logging.getLogger(__name__)

# ============================================================================
# AVAILABILITY CHECKS
# ============================================================================

HTTPX_AVAILABLE = False
PYPPETEER_AVAILABLE = False
SELENIUM_AVAILABLE = False
BEAUTIFULSOUP4_AVAILABLE = False
REQUESTS_AVAILABLE = False
PYPDF_AVAILABLE = False
PYTHON_DOCX_AVAILABLE = False
OPENPYXL_AVAILABLE = False
PYTHON_PPTX_AVAILABLE = False
EBOOKLIB_AVAILABLE = False
PYTESSERACT_AVAILABLE = False
PILLOW_AVAILABLE = False
PYGITHUB_AVAILABLE = False

try:
    import httpx
    HTTPX_AVAILABLE = True
except ImportError:
    pass

try:
    from pyppeteer import launch
    PYPPETEER_AVAILABLE = True
except ImportError:
    pass

try:
    from selenium import webdriver
    from selenium.webdriver.common.by import By
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    SELENIUM_AVAILABLE = True
except ImportError:
    pass

try:
    from bs4 import BeautifulSoup
    BEAUTIFULSOUP4_AVAILABLE = True
except ImportError:
    pass

try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    pass

try:
    from pypdf import PdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    pass

try:
    from docx import Document
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    pass

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    pass

try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    pass

try:
    import ebooklib
    from ebooklib import epub
    EBOOKLIB_AVAILABLE = True
except ImportError:
    pass

try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    pass

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    pass

try:
    from github import Github
    PYGITHUB_AVAILABLE = True
except ImportError:
    pass


# ============================================================================
# DOCUMENT EXTRACTORS
# ============================================================================

def extract_pdf(content_bytes: bytes) -> str:
    """Extract text from PDF using pypdf."""
    if not PYPDF_AVAILABLE:
        return "[PDF extraction unavailable: pypdf not installed]"

    try:
        pdf_file = io.BytesIO(content_bytes)
        reader = PdfReader(pdf_file)
        text_parts = []
        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                text_parts.append(f"[Page {page_num + 1}]\n{text}")
        return "\n\n".join(text_parts) if text_parts else "[PDF extraction: no text found]"
    except Exception as e:
        logger.warning(f"PDF extraction error: {e}")
        return f"[PDF extraction error: {str(e)}]"


def extract_docx(content_bytes: bytes) -> str:
    """Extract text from Word document."""
    if not PYTHON_DOCX_AVAILABLE:
        return "[DOCX extraction unavailable: python-docx not installed]"

    try:
        doc = Document(io.BytesIO(content_bytes))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n".join(paragraphs) if paragraphs else "[DOCX extraction: no text found]"
    except Exception as e:
        logger.warning(f"DOCX extraction error: {e}")
        return f"[DOCX extraction error: {str(e)}]"


def extract_xlsx(content_bytes: bytes) -> str:
    """Extract structured data from Excel as text."""
    if not OPENPYXL_AVAILABLE:
        return "[XLSX extraction unavailable: openpyxl not installed]"

    try:
        wb = load_workbook(io.BytesIO(content_bytes))
        sheet_texts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = f"[Sheet: {sheet_name}]\n"
            rows = []

            for row in ws.iter_rows(values_only=True):
                # Filter out None values and convert to strings
                cells = [str(cell) if cell is not None else "" for cell in row]
                rows.append(" | ".join(cells))

            sheet_text += "\n".join(rows)
            sheet_texts.append(sheet_text)

        return "\n\n".join(sheet_texts) if sheet_texts else "[XLSX extraction: no data found]"
    except Exception as e:
        logger.warning(f"XLSX extraction error: {e}")
        return f"[XLSX extraction error: {str(e)}]"


def extract_pptx(content_bytes: bytes) -> str:
    """Extract text from PowerPoint slides."""
    if not PYTHON_PPTX_AVAILABLE:
        return "[PPTX extraction unavailable: python-pptx not installed]"

    try:
        prs = Presentation(io.BytesIO(content_bytes))
        slide_texts = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = f"[Slide {slide_num + 1}]\n"
            shapes = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes.append(shape.text)

            slide_text += "\n".join(shapes)
            slide_texts.append(slide_text)

        return "\n\n".join(slide_texts) if slide_texts else "[PPTX extraction: no text found]"
    except Exception as e:
        logger.warning(f"PPTX extraction error: {e}")
        return f"[PPTX extraction error: {str(e)}]"


def extract_epub(content_bytes: bytes) -> str:
    """Extract text from EPUB ebook."""
    if not EBOOKLIB_AVAILABLE:
        return "[EPUB extraction unavailable: ebooklib not installed]"

    try:
        book = epub.read_epub(io.BytesIO(content_bytes))
        chapters = []

        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                content = item.get_content()
                if isinstance(content, bytes):
                    content = content.decode('utf-8', errors='ignore')
                # Basic HTML cleanup for EPUB content
                soup = BeautifulSoup(content, 'html.parser') if BEAUTIFULSOUP4_AVAILABLE else None
                if soup:
                    text = soup.get_text(separator='\n', strip=True)
                    chapters.append(text)
                else:
                    chapters.append(content)

        return "\n\n---\n\n".join(chapters) if chapters else "[EPUB extraction: no text found]"
    except Exception as e:
        logger.warning(f"EPUB extraction error: {e}")
        return f"[EPUB extraction error: {str(e)}]"


def extract_image_text(content_bytes: bytes) -> str:
    """OCR text from image using pytesseract."""
    if not PYTESSERACT_AVAILABLE or not PILLOW_AVAILABLE:
        return "[OCR unavailable: pytesseract or Pillow not installed]"

    try:
        image = Image.open(io.BytesIO(content_bytes))
        text = pytesseract.image_to_string(image)
        return text if text.strip() else "[OCR: no text detected in image]"
    except Exception as e:
        logger.warning(f"OCR error: {e}")
        return f"[OCR error: {str(e)}]"


def extract_archive(content_bytes: bytes, url: str) -> str:
    """Extract and process archive contents (ZIP/TAR/GZ)."""
    import zipfile
    import tarfile

    extracted_texts = []

    # Try ZIP first
    try:
        zip_file = zipfile.ZipFile(io.BytesIO(content_bytes))
        for name in zip_file.namelist():
            if name.endswith('/'):
                continue

            try:
                file_bytes = zip_file.read(name)
                text = _process_file_by_extension(name, file_bytes, url)
                if text:
                    extracted_texts.append(f"[File: {name}]\n{text}")
            except Exception as e:
                logger.debug(f"Could not extract {name} from ZIP: {e}")

        if extracted_texts:
            return "\n\n---\n\n".join(extracted_texts)
    except zipfile.BadZipFile:
        pass
    except Exception as e:
        logger.debug(f"ZIP extraction failed: {e}")

    # Try TAR
    try:
        tar_file = tarfile.open(fileobj=io.BytesIO(content_bytes))
        for member in tar_file.getmembers():
            if member.isdir():
                continue

            try:
                f = tar_file.extractfile(member)
                if f:
                    file_bytes = f.read()
                    text = _process_file_by_extension(member.name, file_bytes, url)
                    if text:
                        extracted_texts.append(f"[File: {member.name}]\n{text}")
            except Exception as e:
                logger.debug(f"Could not extract {member.name} from TAR: {e}")

        if extracted_texts:
            return "\n\n---\n\n".join(extracted_texts)
    except tarfile.ReadError:
        pass
    except Exception as e:
        logger.debug(f"TAR extraction failed: {e}")

    return "[Archive extraction: unsupported or corrupted archive]"


def _process_file_by_extension(filename: str, content_bytes: bytes, url: str) -> str:
    """Process a file based on its extension."""
    lower_name = filename.lower()

    if lower_name.endswith('.pdf'):
        return extract_pdf(content_bytes)
    elif lower_name.endswith('.docx'):
        return extract_docx(content_bytes)
    elif lower_name.endswith('.xlsx'):
        return extract_xlsx(content_bytes)
    elif lower_name.endswith('.pptx'):
        return extract_pptx(content_bytes)
    elif lower_name.endswith('.epub'):
        return extract_epub(content_bytes)
    elif lower_name.endswith(('.txt', '.md', '.html', '.xml', '.json', '.csv', '.log')):
        return content_bytes.decode('utf-8', errors='ignore')
    elif lower_name.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
        return extract_image_text(content_bytes)
    else:
        # Try to decode as text
        try:
            return content_bytes.decode('utf-8', errors='ignore')
        except Exception:
            return f"[Unsupported file type: {filename}]"


def extract_html_text(html_content: str) -> str:
    """Extract clean text from HTML using BeautifulSoup."""
    if not BEAUTIFULSOUP4_AVAILABLE:
        # Fallback: basic regex removal of tags
        text = re.sub(r'<[^>]+>', '', html_content)
        return text

    try:
        soup = BeautifulSoup(html_content, 'html.parser')

        # Remove scripts, styles, and metadata
        for element in soup(["script", "style", "meta", "link", "noscript"]):
            element.decompose()

        # Extract text
        text = soup.get_text(separator='\n', strip=True)

        # Clean up excessive whitespace
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return '\n'.join(lines)
    except Exception as e:
        logger.warning(f"HTML text extraction error: {e}")
        # Fallback
        text = re.sub(r'<[^>]+>', '', html_content)
        return text


def fetch_github_file(url: str, github_token: Optional[str] = None) -> Dict[str, Any]:
    """Read file directly from GitHub via API."""
    if not PYGITHUB_AVAILABLE:
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "github_api",
            "error": "PyGithub not available"
        }

    try:
        # Parse GitHub URL
        match = re.match(
            r'https://github\.com/([^/]+)/([^/]+)/blob/([^/]+)/(.*)',
            url
        )
        if not match:
            return {
                "url": url,
                "content": "",
                "raw_html": "",
                "mime_type": "unknown",
                "success": False,
                "method": "github_api",
                "error": "Invalid GitHub URL format"
            }

        owner, repo, branch, file_path = match.groups()

        g = Github(github_token) if github_token else Github()
        repository = g.get_user(owner).get_repo(repo)

        try:
            file_content = repository.get_contents(file_path, ref=branch)
            content = file_content.decoded_content.decode('utf-8', errors='ignore')

            return {
                "url": url,
                "content": content,
                "raw_html": "",
                "mime_type": "text/plain",
                "success": True,
                "method": "github_api",
                "error": None
            }
        except Exception as e:
            logger.warning(f"GitHub API file fetch error: {e}")
            return {
                "url": url,
                "content": "",
                "raw_html": "",
                "mime_type": "unknown",
                "success": False,
                "method": "github_api",
                "error": str(e)
            }

    except Exception as e:
        logger.error(f"GitHub API error: {e}")
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "github_api",
            "error": str(e)
        }


# ============================================================================
# TIER FETCH FUNCTIONS
# ============================================================================

async def _fetch_httpx(url: str, timeout: int = 30) -> Dict[str, Any]:
    """TIER 1: Fast async fetch using httpx."""
    if not HTTPX_AVAILABLE:
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "httpx",
            "error": "httpx not available"
        }

    try:
        async with httpx.AsyncClient(timeout=timeout) as client:
            response = await client.get(url, follow_redirects=True)
            response.raise_for_status()

            content_type = response.headers.get("content-type", "").lower()
            content = response.content

            # Detect MIME type
            mime_type = content_type.split(';')[0] if content_type else None
            if not mime_type or mime_type == "application/octet-stream":
                mime_type, _ = mimetypes.guess_type(url)

            # Handle binary document types
            if mime_type == "application/pdf":
                extracted_text = extract_pdf(content)
                return {
                    "url": url,
                    "content": extracted_text,
                    "raw_html": "",
                    "mime_type": mime_type,
                    "success": True,
                    "method": "httpx+pdf",
                    "error": None
                }

            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                extracted_text = extract_docx(content)
                return {
                    "url": url,
                    "content": extracted_text,
                    "raw_html": "",
                    "mime_type": mime_type,
                    "success": True,
                    "method": "httpx+docx",
                    "error": None
                }

            elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                extracted_text = extract_xlsx(content)
                return {
                    "url": url,
                    "content": extracted_text,
                    "raw_html": "",
                    "mime_type": mime_type,
                    "success": True,
                    "method": "httpx+xlsx",
                    "error": None
                }

            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                extracted_text = extract_pptx(content)
                return {
                    "url": url,
                    "content": extracted_text,
                    "raw_html": "",
                    "mime_type": mime_type,
                    "success": True,
                    "method": "httpx+pptx",
                    "error": None
                }

            elif mime_type == "application/epub+zip":
                extracted_text = extract_epub(content)
                return {
                    "url": url,
                    "content": extracted_text,
                    "raw_html": "",
                    "mime_type": mime_type,
                    "success": True,
                    "method": "httpx+epub",
                    "error": None
                }

            elif mime_type and mime_type.startswith("image/"):
                extracted_text = extract_image_text(content)
                return {
                    "url": url,
                    "content": extracted_text,
                    "raw_html": "",
                    "mime_type": mime_type,
                    "success": True,
                    "method": "httpx+ocr",
                    "error": None
                }

            elif mime_type and mime_type.startswith("application/x-") and mime_type.endswith("zip"):
                extracted_text = extract_archive(content, url)
                return {
                    "url": url,
                    "content": extracted_text,
                    "raw_html": "",
                    "mime_type": mime_type,
                    "success": True,
                    "method": "httpx+archive",
                    "error": None
                }

            # Assume HTML/text
            else:
                try:
                    html_content = content.decode('utf-8', errors='ignore')
                except Exception:
                    html_content = str(content)

                text_content = extract_html_text(html_content)

                return {
                    "url": url,
                    "content": text_content,
                    "raw_html": html_content,
                    "mime_type": mime_type or "text/html",
                    "success": bool(text_content.strip()),
                    "method": "httpx",
                    "error": None
                }

    except Exception as e:
        logger.info(f"httpx fetch failed for {url}: {e}")
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "httpx",
            "error": str(e)
        }


async def _fetch_pyppeteer(url: str, timeout: int = 30) -> Dict[str, Any]:
    """TIER 2: Pure Python Chrome automation via pyppeteer (MAIN JS renderer)."""
    if not PYPPETEER_AVAILABLE:
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "pyppeteer",
            "error": "pyppeteer not available"
        }

    browser = None
    try:
        browser = await launch(headless=True, args=['--no-sandbox', '--disable-setuid-sandbox'])
        page = await browser.newPage()
        await page.goto(url, {"waitUntil": "networkidle2", "timeout": timeout * 1000})

        # Get rendered HTML
        html_content = await page.content()

        # Extract text
        text_content = extract_html_text(html_content)

        # Try to detect MIME type from response
        response = await page.evaluate("() => document.contentType || 'text/html'")
        mime_type = response if response else "text/html"

        await page.close()

        return {
            "url": url,
            "content": text_content,
            "raw_html": html_content,
            "mime_type": mime_type,
            "success": bool(text_content.strip()),
            "method": "pyppeteer",
            "error": None
        }

    except Exception as e:
        logger.info(f"pyppeteer fetch failed for {url}: {e}")
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "pyppeteer",
            "error": str(e)
        }

    finally:
        if browser:
            try:
                await browser.close()
            except Exception:
                pass


def _fetch_selenium(url: str, timeout: int = 30) -> Dict[str, Any]:
    """TIER 3: Selenium WebDriver (battle-tested fallback JS renderer)."""
    if not SELENIUM_AVAILABLE:
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "selenium",
            "error": "selenium not available"
        }

    driver = None
    try:
        # Try to use webdriver-manager if available
        try:
            from webdriver_manager.chrome import ChromeDriverManager
            from selenium.webdriver.chrome.service import Service
            options = webdriver.ChromeOptions()
            options.add_argument('--no-sandbox')
            options.add_argument('--disable-dev-shm-usage')
            options.add_argument('--headless')
            service = Service(ChromeDriverManager().install())
            driver = webdriver.Chrome(service=service, options=options)
        except Exception:
            # Fallback: assume chromedriver in PATH
            options = webdriver.ChromeOptions()
            options.add_argument('--no-sandbox')
            options.add_argument('--disable-dev-shm-usage')
            options.add_argument('--headless')
            driver = webdriver.Chrome(options=options)

        driver.set_page_load_timeout(timeout)
        driver.get(url)

        # Wait for page to load
        try:
            WebDriverWait(driver, min(timeout, 10)).until(
                EC.presence_of_all_elements_located((By.TAG_NAME, "body"))
            )
        except Exception:
            pass

        # Get rendered HTML
        html_content = driver.page_source

        # Extract text
        text_content = extract_html_text(html_content)

        return {
            "url": url,
            "content": text_content,
            "raw_html": html_content,
            "mime_type": "text/html",
            "success": bool(text_content.strip()),
            "method": "selenium",
            "error": None
        }

    except Exception as e:
        logger.info(f"selenium fetch failed for {url}: {e}")
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "selenium",
            "error": str(e)
        }

    finally:
        if driver:
            try:
                driver.quit()
            except Exception:
                pass


def _fetch_requests(url: str, timeout: int = 30) -> Dict[str, Any]:
    """TIER 4: Pure Python synchronous fallback (ALWAYS COMPLETES)."""
    if not REQUESTS_AVAILABLE:
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "requests",
            "error": "requests not available"
        }

    try:
        response = requests.get(url, timeout=timeout, allow_redirects=True)
        response.raise_for_status()

        content_type = response.headers.get("content-type", "").lower()
        content = response.content

        # Detect MIME type
        mime_type = content_type.split(';')[0] if content_type else None
        if not mime_type or mime_type == "application/octet-stream":
            mime_type, _ = mimetypes.guess_type(url)

        # Handle binary document types (same as httpx)
        if mime_type == "application/pdf":
            extracted_text = extract_pdf(content)
            return {
                "url": url,
                "content": extracted_text,
                "raw_html": "",
                "mime_type": mime_type,
                "success": True,
                "method": "requests+pdf",
                "error": None
            }

        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            extracted_text = extract_docx(content)
            return {
                "url": url,
                "content": extracted_text,
                "raw_html": "",
                "mime_type": mime_type,
                "success": True,
                "method": "requests+docx",
                "error": None
            }

        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            extracted_text = extract_xlsx(content)
            return {
                "url": url,
                "content": extracted_text,
                "raw_html": "",
                "mime_type": mime_type,
                "success": True,
                "method": "requests+xlsx",
                "error": None
            }

        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            extracted_text = extract_pptx(content)
            return {
                "url": url,
                "content": extracted_text,
                "raw_html": "",
                "mime_type": mime_type,
                "success": True,
                "method": "requests+pptx",
                "error": None
            }

        elif mime_type == "application/epub+zip":
            extracted_text = extract_epub(content)
            return {
                "url": url,
                "content": extracted_text,
                "raw_html": "",
                "mime_type": mime_type,
                "success": True,
                "method": "requests+epub",
                "error": None
            }

        elif mime_type and mime_type.startswith("image/"):
            extracted_text = extract_image_text(content)
            return {
                "url": url,
                "content": extracted_text,
                "raw_html": "",
                "mime_type": mime_type,
                "success": True,
                "method": "requests+ocr",
                "error": None
            }

        elif mime_type and mime_type.startswith("application/x-") and mime_type.endswith("zip"):
            extracted_text = extract_archive(content, url)
            return {
                "url": url,
                "content": extracted_text,
                "raw_html": "",
                "mime_type": mime_type,
                "success": True,
                "method": "requests+archive",
                "error": None
            }

        # Assume HTML/text
        else:
            try:
                html_content = content.decode('utf-8', errors='ignore')
            except Exception:
                html_content = str(content)

            text_content = extract_html_text(html_content)

            return {
                "url": url,
                "content": text_content,
                "raw_html": html_content,
                "mime_type": mime_type or "text/html",
                "success": bool(text_content.strip()),
                "method": "requests",
                "error": None
            }

    except Exception as e:
        logger.info(f"requests fetch failed for {url}: {e}")
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "requests",
            "error": str(e)
        }


# ============================================================================
# MAIN FETCHER CLASS
# ============================================================================

class WebFetcher:
    """
    Universal web content fetcher.

    Multi-tier fallback chain:
      TIER 1: httpx (fast, pure Python)
      TIER 2: pyppeteer (main JS renderer, pure Python)
      TIER 3: selenium (fallback JS renderer)
      TIER 4: requests (always completes)
    """

    def __init__(self, timeout: int = 30, github_token: Optional[str] = None):
        """
        Initialize the fetcher.

        Args:
            timeout: Request timeout in seconds
            github_token: GitHub API token (optional, for private repos)
        """
        self.timeout = timeout
        self.github_token = github_token or os.getenv("GITHUB_TOKEN")

    async def fetch(self, url: str, force_js: bool = False) -> Dict[str, Any]:
        """
        Fetch and extract content from any web URL.

        Args:
            url: The URL to fetch
            force_js: If True, skip httpx and go straight to pyppeteer

        Returns:
            {
                "url": str,              # Original URL
                "content": str,          # Extracted readable text
                "raw_html": str,         # Original HTML (if applicable)
                "mime_type": str,        # Detected MIME type
                "success": bool,         # True if any content extracted
                "method": str,           # Which tier succeeded
                "error": Optional[str]   # Error message if failed
            }
        """
        # Check if GitHub URL
        if "github.com" in url and "/blob/" in url:
            return fetch_github_file(url, self.github_token)

        # TIER 1: httpx (unless force_js)
        if not force_js:
            logger.debug(f"Attempting TIER 1 (httpx) for {url}")
            result = await _fetch_httpx(url, self.timeout)
            if result["success"]:
                logger.info(f"TIER 1 succeeded: {result['method']}")
                return result
            logger.debug(f"TIER 1 failed: {result['error']}")

        # TIER 2: pyppeteer (MAIN JS renderer)
        logger.debug(f"Attempting TIER 2 (pyppeteer) for {url}")
        result = await _fetch_pyppeteer(url, self.timeout)
        if result["success"]:
            logger.info(f"TIER 2 succeeded: {result['method']}")
            return result
        logger.debug(f"TIER 2 failed: {result['error']}")

        # TIER 3: selenium (FALLBACK JS renderer)
        logger.debug(f"Attempting TIER 3 (selenium) for {url}")
        result = _fetch_selenium(url, self.timeout)
        if result["success"]:
            logger.info(f"TIER 3 succeeded: {result['method']}")
            return result
        logger.debug(f"TIER 3 failed: {result['error']}")

        # TIER 4: requests (BASELINE, always completes)
        logger.debug(f"Attempting TIER 4 (requests) for {url}")
        result = _fetch_requests(url, self.timeout)
        if result["success"]:
            logger.info(f"TIER 4 succeeded: {result['method']}")
            return result
        logger.debug(f"TIER 4 failed: {result['error']}")

        # All tiers failed
        logger.error(f"All fetch tiers failed for {url}")
        return {
            "url": url,
            "content": "",
            "raw_html": "",
            "mime_type": "unknown",
            "success": False,
            "method": "none",
            "error": "All fetch methods failed or unavailable"
        }
